from pptx import Presentation
from pptx.util import Inches, Pt
import main # My file with main programm
from tkinter.messagebox import showinfo, showwarning
import os

def normaliseConfig(parameter):
    value= parameter.split('=')[1]
    try:
        value_inch = Inches(float(value))
        return value_inch
    except:
        return value
def changeStyle(style='Default'):
        print(style)
        with open(f'Styles\\{style}\\config.cfg', 'r') as config:

            DATA = config.readlines()
            
            # For common slides
            global left, top, width, height
            global font_pt, font_family
            global background

            left = normaliseConfig(DATA[0])
            top = normaliseConfig(DATA[1])
            width = normaliseConfig(DATA[2])
            height = normaliseConfig(DATA[3])
            font_pt = normaliseConfig(DATA[4])
            font_family = normaliseConfig(DATA[5])
            background = normaliseConfig(DATA[6])

            # Parameters for first slide
            global left_f, top_f, width_f, height_f 
            global font_pt_f, font_family_f
            global background_f

            left_f = normaliseConfig(DATA[7])
            top_f = normaliseConfig(DATA[8])
            width_f = normaliseConfig(DATA[9])
            height_f = normaliseConfig(DATA[10])
            font_pt_f = normaliseConfig(DATA[11])
            font_family_f = normaliseConfig(DATA[12])
            background_f = normaliseConfig(DATA[13])

def generate(slides_list:dict, file_name='EasyPresentasy', style='Default'):
    presentation = Presentation()
    changeStyle(style)
    savePath = f'C:\\Users\\FiksikFur\\Desktop\\Projects\\Python\\EasyPresentasy\\'
    for i in range(len(slides_list)):
        currentSlide = slides_list[i]
        print(currentSlide.title, currentSlide.subtitle, currentSlide.text, currentSlide.id)
        
        
        if (currentSlide.id == 0):
            if(currentSlide.title != ''):
                first_slide = presentation.slide_layouts[0]
                slide = presentation.slides.add_slide(first_slide)
                
                slide.shapes.title.text = currentSlide.title
                
            if(currentSlide.subtitle != ''):
                slide.placeholders[1].text = currentSlide.subtitle
            
            if(currentSlide.text != ''):
                text_box = slide.shapes.add_textbox(left_f,top_f,width_f,height_f)
                tf = text_box.text_frame

                par = tf.add_paragraph()
                par.text = currentSlide.text

        elif (currentSlide.id > 0):
            other_slide = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(other_slide)

            slide.shapes.title.text = currentSlide.title

            if(currentSlide.subtitle != ''):
                slide.placeholders[1].text = currentSlide.subtitle
            
            if(currentSlide.text != ''):
                text_box = slide.shapes.add_textbox(left,top,width,height)
                tf = text_box.text_frame

                par = tf.add_paragraph()
                par.text = currentSlide.text
                
    for file in os.listdir(savePath):
        if(file == file_name):
            file_name = file_name + str(len(os.listdir(savePath)))
            presentation.save(savePath+file_name+'.pptx')
        else:
            presentation.save(savePath+file_name+'.pptx')

    showinfo("Succesfull", f"File was created with name:\n{file_name}.pptx\nIn directory: My Documents")

if __name__=="__main__":
    showwarning("Can't use", "You can't use module without main programm.")
    quit()